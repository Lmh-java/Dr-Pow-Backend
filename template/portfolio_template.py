import io
from typing import List

from PIL.Image import Image
from pptx import Presentation

from template.base_template import BaseTemplate


# basic template
class PortfolioTemplate(BaseTemplate):

    def __init__(self):
        self.prs = Presentation("template/pptx_templates/Portfolio.pptx")
        self.count = 0

    def create_title_slide(self, title: str) -> None:
        slide_layout = self.prs.slide_layouts[0]  # 0 corresponds to the title slide layout
        slide = self.prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.title
        title_box.text = title

    def create_pic_slide(self, title: str, content: List[str], pic: Image) -> None:
        self.count += 1
        slide_layout = self.prs.slide_layouts[self.count % 2 + 1]
        slide = self.prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.title
        content_box = slide.placeholders[1]
        image_box = slide.placeholders[2]

        title_box.text = title
        for row in content:
            p = content_box.text_frame.add_paragraph()
            run = p.add_run()
            p.level = 0
            run.text = row
        with io.BytesIO() as output:
            pic.save(output, format='PNG')
            image_box.insert_picture(output)

    def get_ppt(self) -> Presentation:
        return self.prs
