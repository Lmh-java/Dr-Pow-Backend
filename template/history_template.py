import io
from typing import List

from PIL.Image import Image
from pptx import Presentation

from template.base_template import BaseTemplate


# basic template
class HistoryTemplate(BaseTemplate):

    def __init__(self):
        self.prs = Presentation("template/pptx_templates/History.pptx")
        self.count = 0

    def create_title_slide(self, title: str) -> None:
        slide_layout = self.prs.slide_layouts[0]  # 0 corresponds to the title slide layout
        slide = self.prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.title
        title_box.text = title

    def create_pic_slide(self, title: str, content: List[str], pic: Image) -> None:
        self.count += 1
        slide_layout = self.prs.slide_layouts[self.count]  # 1 corresponds to the content slide layout
        if self.count + 1 == 4:
            self.count = 0
        slide = self.prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.title
        content_box = slide.placeholders[1]
        image_box = slide.placeholders[10]

        title_box.text = title
        for row in content:
            p = content_box.text_frame.add_paragraph()
            run = p.add_run()
            p.level = 0
            run.text = row
        with io.BytesIO() as output:
            pic.save(output, format='PNG')
            image_box.insert_picture(output)

    def get_ppt(self) -> Presentation:
        return self.prs
